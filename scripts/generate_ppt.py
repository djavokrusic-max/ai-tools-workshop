#!/usr/bin/env python3
"""
PPT 生成脚本 - MoClawny帮你做PPT
用法: python3 generate_ppt.py "<大纲>" <页数> <风格> <输出文件>
"""

import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import uuid

def generate_ppt(outline: str, page_count: int, style: str, output_path: str):
    """生成 PPT 文件"""
    
    # 创建演示文稿
    prs = Presentation()
    
    # 设置幻灯片尺寸 (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 风格配置
    style_configs = {
        '简约': {
            'bg_color': RGBColor(255, 255, 255),
            'title_color': RGBColor(0, 0, 0),
            'content_color': RGBColor(50, 50, 50),
            'accent_color': RGBColor(100, 100, 100),
        },
        '商务': {
            'bg_color': RGBColor(25, 55, 95),
            'title_color': RGBColor(255, 255, 255),
            'content_color': RGBColor(220, 220, 220),
            'accent_color': RGBColor(70, 130, 180),
        },
        '活泼': {
            'bg_color': RGBColor(255, 245, 230),
            'title_color': RGBColor(255, 100, 50),
            'content_color': RGBColor(50, 50, 50),
            'accent_color': RGBColor(255, 180, 100),
        }
    }
    
    config = style_configs.get(style, style_configs['简约'])
    
    # 解析大纲
    lines = [line.strip() for line in outline.split('\n') if line.strip()]
    
    # 生成封面页
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.paragraphs[0].text = "MoClawny帮你做的PPT"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = config['title_color']
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.paragraphs[0].text = f"风格: {style} | 页数: {page_count}"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = config['accent_color']
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 计算每页内容数量
    content_lines = [line for line in lines if line]
    lines_per_page = max(1, len(content_lines) // max(1, page_count - 2))
    
    # 生成内容页
    for i in range(page_count - 2):
        if i >= len(content_lines):
            break
            
        slide = prs.slides.add_slide(slide_layout)
        
        # 页码
        page_num_box = slide.shapes.add_textbox(Inches(0.3), Inches(6.8), Inches(1), Inches(0.5))
        page_num_frame = page_num_box.text_frame
        page_num_frame.paragraphs[0].text = f"{i + 2}"
        page_num_frame.paragraphs[0].font.size = Pt(14)
        page_num_frame.paragraphs[0].font.color.rgb = config['accent_color']
        
        # 内容
        content_start = i * lines_per_page
        content_end = content_start + lines_per_page
        page_content = content_lines[content_start:content_end]
        
        if page_content:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            for j, line in enumerate(page_content):
                if j == 0:
                    p = content_frame.paragraphs[0]
                else:
                    p = content_frame.add_paragraph()
                p.text = line
                p.font.size = Pt(28)
                p.font.color.rgb = config['content_color']
                p.space_after = Pt(12)
    
    # 生成结束页
    slide = prs.slides.add_slide(slide_layout)
    end_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(12.333), Inches(1.5))
    end_frame = end_box.text_frame
    end_frame.paragraphs[0].text = "感谢使用"
    end_frame.paragraphs[0].font.size = Pt(48)
    end_frame.paragraphs[0].font.bold = True
    end_frame.paragraphs[0].font.color.rgb = config['title_color']
    end_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
    sub_frame = sub_box.text_frame
    sub_frame.paragraphs[0].text = "MoClawny AI PPT 工坊 🦞"
    sub_frame.paragraphs[0].font.size = Pt(20)
    sub_frame.paragraphs[0].font.color.rgb = config['accent_color']
    sub_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 保存
    prs.save(output_path)
    return output_path

if __name__ == '__main__':
    if len(sys.argv) < 5:
        print("用法: python3 generate_ppt.py <大纲> <页数> <风格> <输出文件>")
        sys.exit(1)
    
    outline = sys.argv[1]
    page_count = int(sys.argv[2])
    style = sys.argv[3]
    output_path = sys.argv[4]
    
    try:
        result = generate_ppt(outline, page_count, style, output_path)
        print(f"SUCCESS:{result}")
    except Exception as e:
        print(f"ERROR:{str(e)}")
        sys.exit(1)
